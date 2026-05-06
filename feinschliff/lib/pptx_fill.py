"""Fill engine for v2 catalog templates: load .pptx, fill placeholders, emit."""
from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn


def load_template(path: str | Path) -> Presentation:
    prs = Presentation(str(path))
    if len(prs.slides) != 1:
        raise ValueError(
            f"v2 template at {path} has {len(prs.slides)} slides; expected exactly one slide"
        )
    return prs


def _placeholder_by_idx(slide, idx: int):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _clone_paragraph_with_text(template_p, text: str):
    p = etree.fromstring(etree.tostring(template_p))
    for child in list(p):
        local = etree.QName(child.tag).localname
        if local in ("r", "br", "fld"):
            p.remove(child)
    template_r = template_p.find(qn("a:r"))
    rPr_src = None
    if template_r is not None and template_r.find(qn("a:rPr")) is not None:
        rPr_src = etree.fromstring(etree.tostring(template_r.find(qn("a:rPr"))))
    r = etree.Element(qn("a:r"))
    if rPr_src is not None:
        r.append(rPr_src)
    t = etree.SubElement(r, qn("a:t"))
    t.text = text
    insert_index = 0
    for i, child in enumerate(p):
        if etree.QName(child.tag).localname == "pPr":
            insert_index = i + 1
    p.insert(insert_index, r)
    return p


def fill_slot(slide, idx: int, text: str) -> None:
    """Replace placeholder text preserving layout-defined paragraph style."""
    ph = _placeholder_by_idx(slide, idx)
    if ph is None:
        raise KeyError(f"slide has no placeholder idx={idx}")

    # If the slide's own placeholder has a paragraph template (because the
    # template was authored that way), use it. Otherwise fall back to the
    # layout's first paragraph as the style template.
    txBody = ph.text_frame._txBody
    template_p = txBody.find(qn("a:p"))
    if template_p is None:
        layout_ph = None
        for lp in slide.slide_layout.placeholders:
            if lp.placeholder_format.idx == idx:
                layout_ph = lp
                break
        if layout_ph is None:
            ph.text_frame.text = text
            return
        template_p = layout_ph.text_frame._txBody.find(qn("a:p"))

    for p in list(txBody.findall(qn("a:p"))):
        txBody.remove(p)
    for line in text.split("\n"):
        txBody.append(_clone_paragraph_with_text(template_p, line))


def write_filled(prs: Presentation, out_path: str | Path) -> None:
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
