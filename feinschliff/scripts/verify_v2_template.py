"""Visual verification for v2 template ports.

For one (demo deck, layout, v2 template) tuple:
  1. Find the demo slide that uses this layout in the v1 deck.
  2. Read placeholder text from it (idx → text).
  3. Fill the v2 template with the same content.
  4. Render both the v1 demo slide and the v2 filled template to PNG via soffice.
  5. Compare via perceptual hash; pass if distance ≤ threshold.

The script is invoked once per layout. Rendering each takes ~3-6 seconds via
soffice headless. Failed verifications keep the rendered PNGs around for manual
inspection.
"""
from __future__ import annotations

import argparse
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from PIL import Image
import imagehash
from pptx import Presentation


SOFFICE = "/opt/homebrew/bin/soffice"


def _soffice_render(pptx_path: Path, out_dir: Path) -> Path:
    """Convert a single-slide pptx to a single PNG. Returns the PNG path.

    Uses an isolated UserInstallation profile per call: without this, soffice
    instances contend on a shared profile dir and ~30+ sequential invocations
    start failing silently (observed during ferrari batch port).
    """
    out_dir.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="soffice-profile-") as profile:
        env_arg = f"-env:UserInstallation=file://{profile}"
        # PPTX -> PDF
        subprocess.run(
            [SOFFICE, env_arg, "--headless", "--convert-to", "pdf",
             "--outdir", str(out_dir), str(pptx_path)],
            check=True,
            capture_output=True,
        )
    pdf_path = out_dir / pptx_path.with_suffix(".pdf").name
    if not pdf_path.is_file():
        raise RuntimeError(f"soffice produced no PDF for {pptx_path}")
    # PDF -> PNG (96 DPI is plenty for phash; bigger inflates without phash gain)
    subprocess.run(
        ["pdftoppm", "-r", "96", "-png", str(pdf_path), str(out_dir / "page")],
        check=True,
        capture_output=True,
    )
    pngs = sorted(out_dir.glob("page-*.png"))
    if not pngs:
        raise RuntimeError(f"no PNGs produced from {pptx_path}")
    return pngs[0]


def _read_demo_slide_content(demo_deck: Path, layout_name: str) -> tuple[dict[int, str], int]:
    """Return (idx→text mapping for the layout's demo slide, slide index in deck)."""
    prs = Presentation(str(demo_deck))
    for i, slide in enumerate(prs.slides):
        if slide.slide_layout.name == layout_name:
            content: dict[int, str] = {}
            for ph in slide.placeholders:
                idx = ph.placeholder_format.idx
                text = ph.text_frame.text
                if text:
                    content[idx] = text
            return content, i
    raise SystemExit(f"no slide in {demo_deck} uses layout {layout_name!r}")


def _extract_demo_slide_to_pptx(demo_deck: Path, slide_idx: int, out_path: Path) -> None:
    """Save just slide_idx of demo_deck as a single-slide pptx (for baseline rendering)."""
    from pptx.oxml.ns import qn

    shutil.copy(demo_deck, out_path)
    prs = Presentation(str(out_path))
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    sldIdLst_part = prs.slides.part
    for i, sld in enumerate(slides_list):
        if i == slide_idx:
            continue
        rId = sld.get(qn("r:id"))
        sldIdLst_part.drop_rel(rId)
        xml_slides.remove(sld)
    prs.save(out_path)


def _fill_v2_with_content(v2_template: Path, content: dict[int, str], out_path: Path) -> None:
    """Fill a v2 template's placeholders by idx and save."""
    sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
    from lib.pptx_fill import load_template, fill_slot, write_filled

    prs = load_template(str(v2_template))
    slide_idxs = {ph.placeholder_format.idx for ph in prs.slides[0].placeholders}
    # Fill known content first.
    for idx, text in content.items():
        if idx not in slide_idxs:
            continue
        fill_slot(prs.slides[0], idx=idx, text=text)
    # Fill any remaining placeholder with empty string so layout default text doesn't
    # render in the comparison. (The v1 demo slide doesn't have these placeholders at
    # all so they render blank there too.)
    for idx in slide_idxs - set(content):
        try:
            fill_slot(prs.slides[0], idx=idx, text="")
        except Exception:
            pass
    write_filled(prs, out_path)


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--demo-deck", required=True, type=Path)
    ap.add_argument("--layout-name", required=True)
    ap.add_argument("--v2-template", required=True, type=Path)
    ap.add_argument("--out-dir", required=True, type=Path,
                    help="Where to write baseline.png and v2.png for inspection")
    ap.add_argument("--threshold", type=int, default=8,
                    help="Max phash distance (0 = identical, 8 = 'looks different to a human')")
    args = ap.parse_args()

    args.out_dir.mkdir(parents=True, exist_ok=True)

    content, slide_idx = _read_demo_slide_content(args.demo_deck, args.layout_name)

    with tempfile.TemporaryDirectory() as tmp:
        tmp = Path(tmp)
        # 1) Render baseline: single-slide pptx of the demo slide
        baseline_pptx = tmp / "baseline.pptx"
        _extract_demo_slide_to_pptx(args.demo_deck, slide_idx, baseline_pptx)
        baseline_png = _soffice_render(baseline_pptx, tmp / "baseline-render")
        shutil.copy(baseline_png, args.out_dir / "baseline.png")

        # 2) Render v2: fill template with demo content, render
        v2_filled = tmp / "v2-filled.pptx"
        _fill_v2_with_content(args.v2_template, content, v2_filled)
        v2_png = _soffice_render(v2_filled, tmp / "v2-render")
        shutil.copy(v2_png, args.out_dir / "v2.png")

    # 3) phash diff
    h1 = imagehash.phash(Image.open(args.out_dir / "baseline.png"))
    h2 = imagehash.phash(Image.open(args.out_dir / "v2.png"))
    distance = h1 - h2
    verdict = "pass" if distance <= args.threshold else "fail"
    print(f"phash_distance={distance} threshold={args.threshold} verdict={verdict}")
    return 0 if verdict == "pass" else 1


if __name__ == "__main__":
    sys.exit(main())
