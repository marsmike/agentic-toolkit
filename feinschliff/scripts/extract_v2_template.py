"""One-time helper: extract a single layout from the demo deck into a v2 template file.

Two responsibilities beyond the obvious "keep one slide, drop the rest":

1. **drop_rel** — `xml_slides.remove(sld)` only severs the sldIdLst reference. The
   slide XML parts themselves remain in the package, bloating the .pptx by ~4× and
   colliding with the kept slide's renamed slide1.xml on save (LibreOffice rejects
   the duplicate; PowerPoint warns and may auto-repair). `part.drop_rel(rId)`
   removes the underlying part too.

2. **Layout placeholder instantiation** — the demo slide may exercise only a
   subset of the layout's placeholders. fill_slot() addresses placeholders by idx
   on the slide, so any idx defined on the layout but not on the slide silently
   fails to fill. This script copies missing layout placeholder XML into the
   slide so the extracted template is self-contained: any idx in the layout is
   addressable from the slide.
"""
from __future__ import annotations

import argparse
import copy
import hashlib
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn


def _slide_placeholder_idxs(slide) -> set[int]:
    return {ph.placeholder_format.idx for ph in slide.placeholders}


def _instantiate_missing_layout_placeholders(slide, layout) -> int:
    """Copy any layout placeholder XML missing from the slide into the slide. Returns count added."""
    slide_idxs = _slide_placeholder_idxs(slide)
    spTree = slide.shapes._spTree
    added = 0
    for layout_ph in layout.placeholders:
        idx = layout_ph.placeholder_format.idx
        if idx in slide_idxs:
            continue
        sp_clone = copy.deepcopy(layout_ph._element)
        # Strip layout-level positioning hint (a:extLst) if present so the
        # slide-level placeholder inherits geometry from the layout naturally.
        spTree.append(sp_clone)
        added += 1
    return added


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", required=True, help="Built demo deck .pptx")
    ap.add_argument("--layout-name", required=True, help="Slide layout name to extract")
    ap.add_argument("--out", required=True, help="Output single-slide .pptx")
    args = ap.parse_args()

    src = Presentation(args.input)
    target_layout = None
    for sl in src.slide_layouts:
        if sl.name == args.layout_name:
            target_layout = sl
            break
    if target_layout is None:
        layouts = ", ".join(sl.name for sl in src.slide_layouts)
        sys.exit(f"layout {args.layout_name!r} not found. Have: {layouts}")

    demo_slide = None
    for s in src.slides:
        if s.slide_layout.name == args.layout_name:
            demo_slide = s
            break
    if demo_slide is None:
        sys.exit(f"no slide in {args.input} uses layout {args.layout_name!r}")

    out_path = Path(args.out)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    src.save(out_path)
    new = Presentation(out_path)
    keep_idx = None
    for i, s in enumerate(new.slides):
        if s.slide_layout.name == args.layout_name:
            keep_idx = i
            break

    # Materialise every layout placeholder onto the kept slide so fill_slot can
    # address it. (fill_slot resolves placeholders by idx ON THE SLIDE; layout-only
    # placeholders are invisible to it.)
    kept_slide = new.slides[keep_idx]
    added = _instantiate_missing_layout_placeholders(kept_slide, kept_slide.slide_layout)

    # Now drop sldId references AND their underlying parts so the saved .pptx
    # contains only the kept slide.
    xml_slides = new.slides._sldIdLst
    slides_list = list(xml_slides)
    sldIdLst_part = new.slides.part
    dropped = 0
    for i, sld in enumerate(slides_list):
        if i == keep_idx:
            continue
        rId = sld.get(qn("r:id"))
        sldIdLst_part.drop_rel(rId)
        xml_slides.remove(sld)
        dropped += 1

    new.save(out_path)

    sha = hashlib.sha256(out_path.read_bytes()).hexdigest()
    size_kb = out_path.stat().st_size / 1024
    print(f"WROTE {out_path} sha256={sha} size={size_kb:.1f}KB "
          f"(materialised {added} layout-only placeholders, dropped {dropped} orphan slides)")


if __name__ == "__main__":
    main()
