from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from lib.pptx_fill import load_template, fill_slot, write_filled


@pytest.fixture
def minimal_template_pptx(tmp_path):
    """Build a 1-slide pptx with one title placeholder."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    layout = prs.slide_layouts[0]   # built-in title layout
    slide = prs.slides.add_slide(layout)
    out = tmp_path / "minimal.pptx"
    prs.save(out)
    return out


def test_load_template_returns_presentation(minimal_template_pptx):
    prs = load_template(minimal_template_pptx)
    assert len(prs.slides) == 1


def test_load_template_rejects_multi_slide(tmp_path):
    prs = Presentation()
    layout = prs.slide_layouts[0]
    prs.slides.add_slide(layout)
    prs.slides.add_slide(layout)
    multi = tmp_path / "multi.pptx"
    prs.save(multi)
    with pytest.raises(ValueError, match="exactly one slide"):
        load_template(multi)


def test_fill_slot_replaces_placeholder_text(minimal_template_pptx):
    prs = load_template(minimal_template_pptx)
    slide = prs.slides[0]
    fill_slot(slide, idx=0, text="Hello World")
    title_text = slide.placeholders[0].text_frame.text
    assert "Hello World" in title_text


def test_write_filled_emits_valid_pptx(minimal_template_pptx, tmp_path):
    prs = load_template(minimal_template_pptx)
    fill_slot(prs.slides[0], idx=0, text="Out")
    out = tmp_path / "out.pptx"
    write_filled(prs, out)
    assert out.is_file()
    Presentation(out)  # round-trip parse — raises if corrupt
